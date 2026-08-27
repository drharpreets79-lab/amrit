#!/usr/bin/env python3
"""Build the editable AMRIT application data-flow slide.

The diagram is intentionally rendered with native PowerPoint shapes so that
every label, card, arrow, colour and line remains editable.  It is based on the
current desktop capture/import/validation/storage/analysis/sync paths and the
central aggregate query, privacy validation, analytics, dashboard and One
Health control-plane paths in this workspace.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Liberation Sans"


COLORS = {
    "background": "F5F7FA",
    "ink": "172033",
    "muted": "5C687A",
    "white": "FFFFFF",
    "line": "D7DEE8",
    "desktop": "0072B2",       # Okabe-Ito blue
    "desktop_dark": "004E7C",
    "desktop_pale": "EAF4FB",
    "integrator": "D55E00",    # Okabe-Ito vermillion
    "integrator_dark": "8E3F00",
    "integrator_pale": "FFF1E8",
    "exchange": "009E73",      # Okabe-Ito bluish green
    "exchange_dark": "006B4E",
    "exchange_pale": "EAF7F3",
    "amber": "E69F00",         # Okabe-Ito orange
    "amber_dark": "8B5E00",
    "amber_pale": "FFF7E3",
    "purple": "7B61A8",
    "purple_pale": "F3EEFA",
    "slate": "344054",
    "slate_pale": "EEF2F6",
}


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_accessibility(shape, name: str, description: str = "") -> None:
    """Give native shapes stable names and useful alternative descriptions."""
    try:
        props = shape._element.nvSpPr.cNvPr
        props.set("name", name)
        if description:
            props.set("descr", description)
    except AttributeError:
        pass


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str,
    line: str | None = None,
    line_width: float = 1.0,
    radius_name: str = "Shape",
    description: str = "",
):
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS.get(fill, fill))
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(COLORS.get(line, line))
        shape.line.width = Pt(line_width)
    set_accessibility(shape, radius_name, description)
    return shape


def style_text_frame(
    text_frame,
    *,
    margin_left: float = 0.08,
    margin_right: float = 0.08,
    margin_top: float = 0.04,
    margin_bottom: float = 0.04,
    vertical_anchor=MSO_ANCHOR.MIDDLE,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(margin_left)
    text_frame.margin_right = Inches(margin_right)
    text_frame.margin_top = Inches(margin_top)
    text_frame.margin_bottom = Inches(margin_bottom)
    text_frame.vertical_anchor = vertical_anchor


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    font_size: float,
    color: str = "ink",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.MIDDLE,
    margin: float = 0.02,
    name: str = "Text",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    style_text_frame(
        box.text_frame,
        margin_left=margin,
        margin_right=margin,
        margin_top=margin,
        margin_bottom=margin,
        vertical_anchor=vertical_anchor,
    )
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS.get(color, color))
    set_accessibility(box, name, text)
    return box


def add_zone(
    slide,
    x: float,
    w: float,
    *,
    title: str,
    subtitle: str,
    color: str,
    pale: str,
    name: str,
):
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        1.12,
        w,
        5.70,
        fill=pale,
        line=color,
        line_width=1.4,
        radius_name=f"{name} zone",
        description=f"{title}: {subtitle}",
    )
    header = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        1.12,
        w,
        0.62,
        fill=color,
        line=color,
        line_width=1.0,
        radius_name=f"{name} header",
    )
    style_text_frame(header.text_frame, margin_left=0.14, margin_right=0.14)
    title_p = header.text_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_p.space_after = Pt(1)
    title_run = title_p.add_run()
    title_run.text = title
    title_run.font.name = FONT
    title_run.font.size = Pt(17)
    title_run.font.bold = True
    title_run.font.color.rgb = rgb(COLORS["white"])
    sub_p = header.text_frame.add_paragraph()
    sub_p.alignment = PP_ALIGN.CENTER
    sub_p.space_before = Pt(0)
    sub_p.space_after = Pt(0)
    sub_run = sub_p.add_run()
    sub_run.text = subtitle
    sub_run.font.name = FONT
    sub_run.font.size = Pt(8.3)
    sub_run.font.color.rgb = rgb(COLORS["white"])


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    badge: str,
    title: str,
    body: str,
    color: str,
    fill: str = "white",
    name: str,
    title_size: float = 10.7,
    body_size: float = 9.0,
):
    card = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=fill,
        line=color,
        line_width=1.05,
        radius_name=name,
        description=f"{title}. {body}",
    )
    badge_shape = add_shape(
        slide,
        MSO_SHAPE.OVAL,
        x + 0.12,
        y + (h - 0.32) / 2,
        0.32,
        0.32,
        fill=color,
        line=color,
        line_width=0.7,
        radius_name=f"{name} badge",
    )
    style_text_frame(badge_shape.text_frame, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
    badge_p = badge_shape.text_frame.paragraphs[0]
    badge_p.alignment = PP_ALIGN.CENTER
    badge_run = badge_p.add_run()
    badge_run.text = badge
    badge_run.font.name = FONT
    badge_run.font.size = Pt(9.0)
    badge_run.font.bold = True
    badge_run.font.color.rgb = rgb(COLORS["white"])

    tx = x + 0.54
    tw = w - 0.68
    add_text(
        slide,
        tx,
        y + 0.07,
        tw,
        0.24,
        title,
        font_size=title_size,
        color=color,
        bold=True,
        margin=0,
        name=f"{name} title",
    )
    add_text(
        slide,
        tx,
        y + 0.31,
        tw,
        h - 0.36,
        body,
        font_size=body_size,
        color="ink",
        margin=0,
        vertical_anchor=MSO_ANCHOR.TOP,
        name=f"{name} detail",
    )
    return card


def add_flow_arrow(
    slide,
    x: float,
    y: float,
    *,
    direction: str,
    color: str,
    name: str,
):
    kind = {
        "down": MSO_SHAPE.DOWN_ARROW,
        "up_down": MSO_SHAPE.UP_DOWN_ARROW,
    }[direction]
    arrow_w, arrow_h = ((0.14, 0.18) if direction == "up_down" else (0.22, 0.16))
    return add_shape(
        slide,
        kind,
        x,
        y,
        arrow_w,
        arrow_h,
        fill=color,
        line=color,
        line_width=0.5,
        radius_name=name,
    )


def add_exchange_arrow(
    slide,
    y: float,
    *,
    direction: str,
    label: str,
    detail: str,
    color: str,
    name: str,
):
    shape_type = MSO_SHAPE.LEFT_ARROW if direction == "left" else MSO_SHAPE.RIGHT_ARROW
    arrow = add_shape(
        slide,
        shape_type,
        5.29,
        y,
        2.73,
        0.78,
        fill=color,
        line=color,
        line_width=1.0,
        radius_name=name,
        description=f"{label}. {detail}",
    )
    style_text_frame(
        arrow.text_frame,
        margin_left=0.38 if direction == "left" else 0.14,
        margin_right=0.14 if direction == "left" else 0.38,
        margin_top=0.07,
        margin_bottom=0.04,
    )
    p = arrow.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = label
    run.font.name = FONT
    run.font.size = Pt(10.8)
    run.font.bold = True
    run.font.color.rgb = rgb(COLORS["white"])
    p2 = arrow.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(0)
    p2.space_after = Pt(0)
    run2 = p2.add_run()
    run2.text = detail
    run2.font.name = FONT
    run2.font.size = Pt(8.5)
    run2.font.color.rgb = rgb(COLORS["white"])
    return arrow


def build_presentation(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "AMRIT application data flow"
    prs.core_properties.subject = "Editable data-flow diagram for AMRIT Desktop and AMRIT Centerlized Integrator"
    prs.core_properties.author = "AMRIT"
    prs.core_properties.keywords = "AMRIT, desktop, centralized integrator, aggregate data, FHIR, One Health"
    prs.core_properties.comments = "All diagram elements are editable native PowerPoint shapes."

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["background"])

    # Title area
    add_text(
        slide,
        0.48,
        0.22,
        12.35,
        0.42,
        "AMRIT application data flow",
        font_size=28,
        color="ink",
        bold=True,
        margin=0,
        name="Diagram title",
    )
    add_text(
        slide,
        0.49,
        0.68,
        12.2,
        0.24,
        "Full application • simplified operational view • editable native shapes",
        font_size=10.5,
        color="muted",
        margin=0,
        name="Diagram subtitle",
    )
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        0.48,
        0.98,
        12.37,
        0.025,
        fill="line",
        radius_name="Title divider",
    )

    # The three requested zones.
    add_zone(
        slide,
        0.45,
        4.34,
        title="AMRIT Desktop",
        subtitle="offline laboratory capture, validation and local analysis",
        color="desktop",
        pale="desktop_pale",
        name="Desktop",
    )
    add_zone(
        slide,
        4.95,
        3.42,
        title="SECURE DATA EXCHANGE",
        subtitle="bidirectional control flow • aggregate data outward",
        color="slate",
        pale="slate_pale",
        name="Exchange",
    )
    add_zone(
        slide,
        8.53,
        4.35,
        title="AMRIT Centerlized Integrator",
        subtitle="national aggregate, analytics and governance layer",
        color="integrator",
        pale="integrator_pale",
        name="Integrator",
    )

    # AMRIT Desktop internal flow.
    card_x = 0.67
    card_w = 3.90
    add_card(
        slide,
        card_x,
        1.88,
        card_w,
        0.62,
        badge="1",
        title="DATA SOURCES",
        body="Manual entry • Excel / CSV • LIS / WHONET",
        color="desktop",
        fill="white",
        name="Desktop data sources",
        body_size=9.2,
    )
    add_flow_arrow(slide, 2.52, 2.51, direction="down", color="desktop", name="Desktop flow 1 to 2")
    add_card(
        slide,
        card_x,
        2.69,
        card_w,
        0.66,
        badge="2",
        title="CAPTURE & MAP",
        body="Forms • protected template • reusable column profiles",
        color="exchange",
        fill="white",
        name="Desktop capture and mapping",
        body_size=9.0,
    )
    add_flow_arrow(slide, 2.52, 3.37, direction="down", color="exchange", name="Desktop flow 2 to 3")
    add_card(
        slide,
        card_x,
        3.54,
        card_w,
        0.82,
        badge="3",
        title="VALIDATE WITH CUSTOM MASTERS & PANELS",
        body="Required fields • organism / sample codes • panel match • AST and duplicate checks",
        color="amber_dark",
        fill="amber_pale",
        name="Desktop validation and panel matching",
        title_size=10.2,
        body_size=8.8,
    )
    add_flow_arrow(slide, 2.52, 4.38, direction="down", color="amber_dark", name="Desktop flow 3 to 4")
    add_card(
        slide,
        card_x,
        4.54,
        card_w,
        0.74,
        badge="4",
        title="LOCAL SQLITE DATA LAYER",
        body="Isolates • masters • import profiles • audit and durable outbox",
        color="purple",
        fill="purple_pale",
        name="Desktop local data layer",
        body_size=8.9,
    )
    add_flow_arrow(slide, 2.52, 5.30, direction="down", color="purple", name="Desktop flow 4 to 5")
    add_card(
        slide,
        card_x,
        5.46,
        card_w,
        1.02,
        badge="5",
        title="ANALYZE, EXPORT & SYNCHRONIZE",
        body="Trends and expert alerts • WHONET / HL7 / FHIR exports • de-identified aggregates",
        color="desktop_dark",
        fill="white",
        name="Desktop analytics exports and sync",
        title_size=10.4,
        body_size=9.0,
    )
    add_text(
        slide,
        0.82,
        6.55,
        3.60,
        0.16,
        "Patient-level records remain on the device",
        font_size=8.5,
        color="desktop_dark",
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
        name="Desktop privacy note",
    )

    # Exchange between both application blocks.
    add_text(
        slide,
        5.19,
        1.89,
        2.94,
        0.18,
        "INTEGRATOR → DESKTOP",
        font_size=8.3,
        color="muted",
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
        name="Request direction label",
    )
    add_exchange_arrow(
        slide,
        2.14,
        direction="left",
        label="Queries & live refresh",
        detail="long-poll • WebSocket",
        color="desktop",
        name="Queries from integrator to desktop",
    )
    add_text(
        slide,
        5.19,
        3.15,
        2.94,
        0.18,
        "DESKTOP → INTEGRATOR",
        font_size=8.3,
        color="muted",
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
        name="Aggregate direction label",
    )
    add_exchange_arrow(
        slide,
        3.40,
        direction="right",
        label="Heartbeat & status",
        detail="aggregate JSON / FHIR",
        color="integrator",
        name="Aggregates from desktop to integrator",
    )
    add_text(
        slide,
        5.19,
        4.54,
        2.94,
        0.18,
        "DESKTOP → INTEGRATOR",
        font_size=8.3,
        color="muted",
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
        name="Product direction label",
    )
    add_exchange_arrow(
        slide,
        4.79,
        direction="right",
        label="One Health products",
        detail="aggregate + lineage / quality",
        color="exchange",
        name="One Health products from desktop to integrator",
    )
    security = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        5.28,
        5.84,
        2.75,
        0.64,
        fill="exchange_dark",
        line="exchange_dark",
        line_width=1.0,
        radius_name="Exchange security controls",
        description="HTTPS or WSS, bearer and site token, privacy validation and audit.",
    )
    style_text_frame(security.text_frame, margin_left=0.10, margin_right=0.10)
    p = security.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "HTTPS / WSS  •  bearer + site token"
    r.font.name = FONT
    r.font.size = Pt(9.5)
    r.font.bold = True
    r.font.color.rgb = rgb(COLORS["white"])
    p2 = security.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "privacy validation • retry • audit"
    r2.font.name = FONT
    r2.font.size = Pt(8.3)
    r2.font.color.rgb = rgb(COLORS["white"])
    add_text(
        slide,
        5.28,
        6.55,
        2.75,
        0.16,
        "No patient identifiers cross this boundary",
        font_size=8.5,
        color="exchange_dark",
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
        name="Exchange privacy note",
    )

    # AMRIT Centerlized Integrator internal flow.  Bidirectional arrows show
    # that portal queries travel down while aggregate results travel upward.
    right_x = 8.75
    right_w = 3.91
    add_card(
        slide,
        right_x,
        1.88,
        right_w,
        0.62,
        badge="5",
        title="STAKEHOLDERS & SYSTEMS",
        body="National • state • district • hospital • analyst / public",
        color="integrator_dark",
        fill="white",
        name="Integrator stakeholders",
        body_size=8.9,
    )
    add_flow_arrow(slide, 10.60, 2.52, direction="up_down", color="integrator", name="Integrator flow 4 and 5")
    add_card(
        slide,
        right_x,
        2.69,
        right_w,
        0.76,
        badge="4",
        title="ROLE PORTALS & APIs",
        body="Dashboards • action inbox • One Health workbench • FHIR / JSON / CSV",
        color="integrator",
        fill="white",
        name="Integrator portals and APIs",
        body_size=8.8,
    )
    add_flow_arrow(slide, 10.60, 3.47, direction="up_down", color="integrator", name="Integrator flow 3 and 4")
    add_card(
        slide,
        right_x,
        3.64,
        right_w,
        0.82,
        badge="3",
        title="METRICS & GOVERNANCE",
        body="Aggregate roll-up + k-anonymity • KPI snapshots • alerts • action plans",
        color="exchange_dark",
        fill="exchange_pale",
        name="Integrator analytics and governance",
        body_size=8.8,
    )
    add_flow_arrow(slide, 10.60, 4.48, direction="up_down", color="exchange", name="Integrator flow 2 and 3")
    add_card(
        slide,
        right_x,
        4.65,
        right_w,
        0.76,
        badge="2",
        title="AGGREGATE DATA LAYER",
        body="Sites • queries / results • snapshots • data products • audit / lineage",
        color="purple",
        fill="purple_pale",
        name="Integrator aggregate data layer",
        body_size=8.8,
    )
    add_flow_arrow(slide, 10.60, 5.43, direction="up_down", color="purple", name="Integrator flow 1 and 2")
    add_card(
        slide,
        right_x,
        5.60,
        right_w,
        0.88,
        badge="1",
        title="SITE EXCHANGE GATEWAY",
        body="Site authentication • query dispatch • PII guard • aggregate FHIR validation",
        color="integrator_dark",
        fill="white",
        name="Integrator exchange gateway",
        body_size=8.8,
    )
    add_text(
        slide,
        8.88,
        6.55,
        3.65,
        0.16,
        "Stores aggregate and control-plane data only",
        font_size=8.5,
        color="integrator_dark",
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
        name="Integrator privacy note",
    )

    # Footer privacy statement is the central architectural invariant.
    footer = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.45,
        6.96,
        12.43,
        0.36,
        fill="ink",
        line="ink",
        line_width=0.7,
        radius_name="Privacy boundary footer",
        description=(
            "Patient-level records stay in AMRIT Desktop. AMRIT Centerlized Integrator accepts "
            "aggregate de-identified responses and governed aggregate products."
        ),
    )
    style_text_frame(footer.text_frame, margin_left=0.16, margin_right=0.16, margin_top=0.02, margin_bottom=0.02)
    fp = footer.text_frame.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = (
        "PRIVACY BOUNDARY  |  Patient-level records stay in AMRIT Desktop; "
        "AMRIT Centerlized Integrator accepts aggregate, de-identified data only."
    )
    fr.font.name = FONT
    fr.font.size = Pt(10.0)
    fr.font.bold = True
    fr.font.color.rgb = rgb(COLORS["white"])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output",
        type=Path,
        default=Path(__file__).resolve().parent / "AMRIT_Application_Data_Flow.pptx",
    )
    args = parser.parse_args()
    build_presentation(args.output)
    print(args.output)


if __name__ == "__main__":
    main()
